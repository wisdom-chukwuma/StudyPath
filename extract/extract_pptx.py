"""Dump slide text from a .pptx as JSON, for turning lecture slides into lesson content.

Usage: python extract_pptx.py "path/to/deck.pptx" > out.json
"""
import sys
import json
from pptx import Presentation


def extract(path):
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    texts.append(text)
        slides.append({"slide": i, "text": texts})
    return slides


if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("Usage: python extract_pptx.py <path-to-pptx>", file=sys.stderr)
        sys.exit(1)
    print(json.dumps(extract(sys.argv[1]), indent=2, ensure_ascii=False))
